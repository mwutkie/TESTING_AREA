import collections 
import collections.abc
from pptx import Presentation
from pptx.chart.data import CategoryChartData

prs = Presentation('xd.pptx')
x=prs.slides[0].shapes[0].chart


chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
x.replace_data(chart_data)
prs.save('hehexd.pptx')

